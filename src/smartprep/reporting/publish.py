"""Publishing: PDF, PowerPoint and notebook, from the same specs.

Every format renders the *same* ``ChartSpec`` objects and the same EDA numbers
that the HTML report and the Studio use. That is the point of the spec layer --
a figure in the PDF and the same figure on screen cannot disagree, because
neither is authored separately.

Each publisher degrades honestly: a missing optional dependency produces a
message naming the install command, never a partial file.
"""

from __future__ import annotations

import pathlib
from dataclasses import dataclass, field
from typing import TYPE_CHECKING, Any

from ..viz.renderers import _require, to_matplotlib
from ..viz.spec import ChartSpec
from ..viz.svg import LIGHT, Theme

if TYPE_CHECKING:  # pragma: no cover
    from ..prepare import PreparationResult

__all__ = ["Slide", "Deck", "to_pdf", "to_pptx", "to_notebook", "publish"]


@dataclass
class Slide:
    """One page or slide: a heading, some prose, some figures, some rows.

    Deliberately format-agnostic. The same structure becomes a PDF page, a
    PowerPoint slide or a notebook cell.
    """

    title: str
    body: str = ""
    charts: list[ChartSpec] = field(default_factory=list)
    table: tuple[list[str], list[list[str]]] | None = None
    notes: str = ""


@dataclass
class Deck:
    """An ordered set of slides, built once and published anywhere."""

    title: str
    subtitle: str = ""
    slides: list[Slide] = field(default_factory=list)

    def add(self, slide: Slide) -> Deck:
        self.slides.append(slide)
        return self

    def __len__(self) -> int:
        return len(self.slides)


def build_deck(result: PreparationResult, *, title: str = "Data Preparation") -> Deck:
    """Assemble the standard report as slides.

    The running order is the argument the report makes: what the data was, what
    was found, what changed, and -- never last and never omitted -- what was
    deliberately left alone.
    """
    from ..eda import compare_profiles, missingness, profile
    from ..viz import (
        before_after_chart,
        health_chart,
        issue_chart,
        kpi_chart,
        missingness_chart,
        stage_chart,
    )

    before, after = result.health_before, result.health_after
    comparison = compare_profiles(profile(result.raw_df), profile(result.clean_df))

    deck = Deck(title=title, subtitle=f"Status: {result.status.value}")

    deck.add(
        Slide(
            title="Summary",
            body=(
                f"Status {result.status.value}. Data health moved from "
                f"{before.overall:.0f} to {after.overall:.0f} across "
                f"{len(result.audit.applied)} operations, changing "
                f"{result.cells_changed} cells. "
                f"{len(result.review_queue)} findings still need a decision."
            ),
            charts=[
                c
                for c in (
                    kpi_chart(
                        [
                            ("Data health", f"{before.overall:.0f} to {after.overall:.0f}"),
                            ("Operations", str(len(result.audit.applied))),
                            ("Cells changed", str(result.cells_changed)),
                            ("Still open", str(len(result.review_queue))),
                        ]
                    ),
                    health_chart(before, after),
                )
                if c is not None
            ],
        )
    )

    deck.add(
        Slide(
            title="Findings",
            body=(
                "Every finding is classified by whether it can be repaired without "
                "asking, and if not, by whose decision it needs."
            ),
            charts=[
                c
                for c in (
                    issue_chart(result.before_scan.issues),
                    missingness_chart(missingness(result.raw_df)),
                )
                if c is not None
            ],
            table=(
                ["Issue", "Class", "Rows", "Summary"],
                [
                    [i.id, i.repair_class.name, str(i.affected_row_count), i.evidence.summary[:80]]
                    for i in result.before_scan.issues[:14]
                ],
            ),
        )
    )

    stages = [
        ("raw", before.overall),
        ("after safe repairs", after.overall),
    ]
    deck.add(
        Slide(
            title="What changed",
            body=(
                "A repair that improves completeness can still move the mean and "
                "shrink the variance. Both are shown."
            ),
            charts=[
                c
                for c in (
                    before_after_chart(comparison),
                    stage_chart(
                        stages,
                        title="Data health by stage",
                        reason="ordered steps of one process, so the axis carries meaning",
                    ),
                )
                if c is not None
            ],
            table=(
                ["Where", "Red flag"],
                [[w, t] for w, t in comparison.red_flags[:12]] or [["—", "none"]],
            ),
        )
    )

    open_rows = [
        [
            i.id,
            i.repair_class.name,
            str(i.affected_row_count),
            "; ".join(i.abstention_reasons)[:90] or i.evidence.summary[:90],
        ]
        for i in result.review_queue[:16]
    ]
    deck.add(
        Slide(
            title="What auto mode did NOT do",
            body=(
                "Mandatory. clean_df is not a verified dataset: these findings "
                "remain open, with the reason each was left alone."
            ),
            table=(
                ["Issue", "Class", "Rows", "Why"],
                open_rows or [["—", "—", "—", "nothing open"]],
            ),
            notes="Never place this section in an appendix.",
        )
    )

    return deck


# --------------------------------------------------------------------------
# PDF
# --------------------------------------------------------------------------


#: Printed in every PDF. The distinctions the library is built on are the ones
#: a reader is most likely to collapse, so they are stated where the numbers
#: are, not only in the documentation.
_METHODOLOGY = """Scan coverage is not data health. Coverage counts the checks that ran;
health describes what they found. A dataset can be fully scanned and
badly broken, or barely scanned and sound.

Detection confidence is not repair confidence. A defect can be certain
while its correct repair is unknown -- 31/02/2025 is definitely invalid
and there is no way to know what date was meant. Only repair confidence
reaches the autonomy ladder, so a certain defect with an uncertain fix
is escalated, never guessed at.

clean_df is not verified_df. Automatic mode repairs only what it can
justify; everything else is reported and left alone. A dataset becomes
verified when a reviewer finalizes it, with any remaining findings
waived on the record.

Cleaning is not preprocessing. Nothing that requires a target, a split
or a modelling decision runs automatically.

Charts are rendered from the same specifications as the screen report.
A figure here and the same figure on screen are drawn from one set of
numbers and cannot disagree.

Sampled charts say so in their caption. A chart that does not say it was
sampled was drawn from the full data."""


#: A4 portrait, in inches. Landscape is the same page turned, used for tables
#: too wide to read down the page.
_PORTRAIT = (8.27, 11.69)
_LANDSCAPE = (11.69, 8.27)


def _pages_for(slide: Slide) -> int:
    """How many PDF pages a slide occupies.

    Deterministic, so the table of contents can be numbered in one pass
    rather than guessed at.
    """
    return 1 + len(slide.charts) + (1 if slide.table else 0)


def _chrome(figure: Any, deck: Deck, number: int, total: int, theme: Theme) -> None:
    """Running header and footer.

    A printed page that has become separated from its report should still say
    what it belongs to and where it sat.
    """
    figure.text(0.08, 0.975, deck.title, fontsize=7, color=theme.muted)
    figure.text(0.92, 0.975, deck.subtitle, fontsize=7, color=theme.muted, ha="right")
    figure.text(0.5, 0.02, f"{number} of {total}", fontsize=7, color=theme.muted, ha="center")


def to_pdf(
    result: PreparationResult | Deck,
    path: str,
    *,
    theme: Theme = LIGHT,
    title: str = "Data Preparation",
) -> str:
    """Write a multi-page PDF with publication-quality static figures.

    Cover, contents, body, then a methodology appendix. Every page carries a
    running header and a page number, and every figure carries a caption
    naming what it shows and why it was drawn -- a figure a reader cannot
    cite is a figure they cannot argue with.

    Charts are rendered from the same ``ChartSpec`` objects as the screen
    report, with interaction lowered on the spec rather than reinterpreted by
    the renderer: paper cannot hover.
    """
    _require("matplotlib", "viz")
    import matplotlib.pyplot as plt
    from matplotlib.backends.backend_pdf import PdfPages

    deck = result if isinstance(result, Deck) else build_deck(result, title=title)
    target = pathlib.Path(path)

    # Cover + contents + body + methodology appendix.
    spans = [_pages_for(slide) for slide in deck.slides]
    total = 2 + sum(spans) + 1
    contents: list[tuple[str, int]] = []
    cursor = 3
    for slide, span in zip(deck.slides, spans, strict=True):
        contents.append((slide.title, cursor))
        cursor += span
    contents.append(("Methodology and caveats", cursor))

    number = 0

    def page(landscape: bool = False) -> Any:
        nonlocal number
        number += 1
        figure = plt.figure(figsize=_LANDSCAPE if landscape else _PORTRAIT)
        figure.patch.set_facecolor(theme.background)
        return figure

    with PdfPages(target) as pdf:
        cover = page()
        cover.text(0.08, 0.86, deck.title, fontsize=26, fontweight="bold", color=theme.foreground)
        cover.text(0.08, 0.81, deck.subtitle, fontsize=12, color=theme.muted)
        cover.text(
            0.08,
            0.06,
            "Scan coverage measures checks executed, not data correctness. "
            "clean_df is not a verified dataset.",
            fontsize=8,
            color=theme.muted,
            wrap=True,
        )
        _chrome(cover, deck, number, total, theme)
        pdf.savefig(cover)
        plt.close(cover)

        toc = page()
        toc.text(0.08, 0.92, "Contents", fontsize=17, fontweight="bold", color=theme.foreground)
        for index, (heading, at) in enumerate(contents):
            y = 0.86 - index * 0.033
            if y < 0.08:
                break
            toc.text(0.08, y, heading, fontsize=10, color=theme.foreground)
            toc.text(0.92, y, str(at), fontsize=10, color=theme.muted, ha="right")
        _chrome(toc, deck, number, total, theme)
        pdf.savefig(toc)
        plt.close(toc)

        for slide in deck.slides:
            text_page = page()
            text_page.text(
                0.08, 0.92, slide.title, fontsize=17, fontweight="bold", color=theme.foreground
            )
            if slide.body:
                text_page.text(
                    0.08, 0.88, slide.body, fontsize=9, color=theme.muted, wrap=True, va="top"
                )
            _chrome(text_page, deck, number, total, theme)
            pdf.savefig(text_page)
            plt.close(text_page)

            for position, chart in enumerate(slide.charts, start=1):
                number += 1
                # as_static() lowers the ceiling on the spec, so the renderer
                # never has to decide for itself what print means.
                figure = to_matplotlib(chart.as_static(), theme)
                caption = f"Figure {number}.{position} -- {chart.title or 'chart'}"
                if chart.rationale:
                    caption += f". {chart.rationale}"
                if chart.is_sampled:
                    caption += f"  [{chart.fidelity.value}: {chart.fidelity_note}]"
                figure.text(0.02, -0.04, caption[:230], fontsize=7, color=theme.muted, wrap=True)
                _chrome(figure, deck, number, total, theme)
                pdf.savefig(figure, bbox_inches="tight")
                plt.close(figure)

            if slide.table:
                headers, rows = slide.table
                table_page = page(landscape=True)
                axes = table_page.add_subplot(111)
                axes.set_axis_off()
                shown = rows[:24]
                axes.set_title(
                    f"{slide.title} -- detail"
                    + (f" (first 24 of {len(rows)} rows)" if len(rows) > 24 else ""),
                    fontsize=12,
                    loc="left",
                )
                rendered = axes.table(
                    cellText=shown,
                    colLabels=headers,
                    loc="upper center",
                    cellLoc="left",
                    colLoc="left",
                )
                rendered.auto_set_font_size(False)
                rendered.set_fontsize(7)
                rendered.scale(1, 1.35)
                if len(rows) > len(shown):
                    # Table continuation: say what was cut, rather than let a
                    # reader assume the page held everything.
                    table_page.text(
                        0.08,
                        0.06,
                        f"{len(rows) - len(shown)} further rows omitted for print; "
                        "the full table is in the HTML report and in Python.",
                        fontsize=8,
                        color=theme.muted,
                    )
                _chrome(table_page, deck, number, total, theme)
                pdf.savefig(table_page, bbox_inches="tight")
                plt.close(table_page)

        appendix = page()
        appendix.text(
            0.08,
            0.92,
            "Methodology and caveats",
            fontsize=17,
            fontweight="bold",
            color=theme.foreground,
        )
        appendix.text(0.08, 0.87, _METHODOLOGY, fontsize=8.5, color=theme.muted, va="top")
        _chrome(appendix, deck, number, total, theme)
        pdf.savefig(appendix)
        plt.close(appendix)

        info = pdf.infodict()
        info["Title"] = deck.title
        info["Subject"] = deck.subtitle
        info["Creator"] = "SmartPrep"

    return str(target.resolve())


# --------------------------------------------------------------------------
# PowerPoint
# --------------------------------------------------------------------------


def to_pptx(
    result: PreparationResult | Deck,
    path: str,
    *,
    theme: Theme = LIGHT,
    title: str = "Data Preparation",
) -> str:
    """Write a PowerPoint deck, with charts rendered as embedded images."""
    _require("pptx", "pptx")
    import tempfile

    import matplotlib.pyplot as plt
    from pptx import Presentation
    from pptx.util import Inches, Pt

    deck = result if isinstance(result, Deck) else build_deck(result, title=title)
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    title_layout = presentation.slide_layouts[0]

    cover = presentation.slides.add_slide(title_layout)
    cover.shapes.title.text = deck.title
    cover.placeholders[1].text = deck.subtitle

    scratch = pathlib.Path(tempfile.mkdtemp())
    try:
        for slide_spec in deck.slides:
            slide = presentation.slides.add_slide(blank)
            heading = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
            heading.text_frame.text = slide_spec.title
            heading.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
            heading.text_frame.paragraphs[0].runs[0].font.bold = True

            if slide_spec.body:
                body = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(1.0))
                body.text_frame.word_wrap = True
                body.text_frame.text = slide_spec.body
                for paragraph in body.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(12)

            for index, chart in enumerate(slide_spec.charts[:2]):
                # Slides are print too: the spec's ceiling is lowered here
                # rather than the renderer deciding on its own.
                figure = to_matplotlib(chart.as_static(), theme)
                image = scratch / f"{id(chart)}_{index}.png"
                figure.savefig(
                    image, dpi=170, bbox_inches="tight", facecolor=figure.get_facecolor()
                )
                plt.close(figure)
                slide.shapes.add_picture(
                    str(image), Inches(0.5 + index * 4.7), Inches(2.1), width=Inches(4.4)
                )

            if slide_spec.table:
                headers, rows = slide_spec.table
                capped = rows[:8]
                table_shape = slide.shapes.add_table(
                    len(capped) + 1,
                    len(headers),
                    Inches(0.5),
                    Inches(2.1 if not slide_spec.charts else 5.6),
                    Inches(9),
                    Inches(0.35 * (len(capped) + 1)),
                ).table
                for column, header in enumerate(headers):
                    table_shape.cell(0, column).text = str(header)
                for row_index, row in enumerate(capped, start=1):
                    for column, value in enumerate(row):
                        table_shape.cell(row_index, column).text = str(value)[:70]

            if slide_spec.notes:
                slide.notes_slide.notes_text_frame.text = slide_spec.notes

        presentation.save(str(pathlib.Path(path)))
    finally:
        import shutil

        shutil.rmtree(scratch, ignore_errors=True)

    return str(pathlib.Path(path).resolve())


# --------------------------------------------------------------------------
# Notebook
# --------------------------------------------------------------------------


def to_notebook(result: PreparationResult, path: str, **context: Any) -> str:
    """Write a runnable notebook that reproduces the analysis.

    Not a transcript of what happened -- executable code that redoes it. A
    report the reader can run is a report they can disagree with.
    """
    import json

    arguments = ", ".join(f"{k}={v!r}" for k, v in context.items())
    call = f"sp.auto_prepare(df{', ' + arguments if arguments else ''})"

    def code(source: str) -> dict[str, Any]:
        return {
            "cell_type": "code",
            "execution_count": None,
            "metadata": {},
            "outputs": [],
            "source": source.strip().splitlines(keepends=True),
        }

    def markdown(source: str) -> dict[str, Any]:
        return {
            "cell_type": "markdown",
            "metadata": {},
            "source": source.strip().splitlines(keepends=True),
        }

    open_list = (
        "\n".join(
            f"- `{i.id}` [{i.repair_class.name}] {i.evidence.summary}"
            for i in result.review_queue[:12]
        )
        or "- none"
    )

    notebook = {
        "cells": [
            markdown(
                f"# Data preparation\n\n"
                f"Status **{result.status.value}**. Health "
                f"{result.health_before.overall:.0f} to "
                f"{result.health_after.overall:.0f}.\n\n"
                "Every cell below is runnable: this reproduces the analysis rather "
                "than describing it."
            ),
            code("import pandas as pd\nimport smartprep as sp"),
            code('df = pd.read_excel("your_data.xlsx", dtype=object)  # your source here'),
            markdown("## Scan\n\nDiagnosis only. Nothing is modified."),
            code("scan = sp.scan(df)\nprint(scan.summary())"),
            markdown("## Profile and EDA"),
            code(
                "profile = sp.profile(df)\n"
                "print(profile.summary())\n"
                "print(sp.associations(df, profile).summary())\n"
                "print(sp.missingness(df).summary())"
            ),
            markdown("## Safe automatic preparation"),
            code(f"result = {call}\nprint(result.summary())"),
            markdown(
                "## What auto mode did not do\n\n"
                "`clean_df` is not a verified dataset. These remain open:\n\n"
                f"{open_list}"
            ),
            code("print(result.what_auto_mode_did_not_do())"),
            markdown("## Guided review\n\nResolve what automatic mode declined to decide."),
            code(
                "session = result.open_guided()\n"
                "print(session.summary())\n"
                "# q = session.next_question(); print(q.render())\n"
                "# session.answer(q.issue_id, 'use_recommendation')\n"
                "# final = session.finish()"
            ),
            markdown("## Before and after"),
            code("comparison = result.compare_profiles()\nprint(comparison.summary())"),
            markdown("## Reports"),
            code('result.export_report("report.html")\nsp.studio(result)   # renders inline'),
        ],
        "metadata": {
            "kernelspec": {"display_name": "Python 3", "language": "python", "name": "python3"},
            "language_info": {"name": "python"},
        },
        "nbformat": 4,
        "nbformat_minor": 5,
    }

    target = pathlib.Path(path)
    target.write_text(json.dumps(notebook, indent=1, ensure_ascii=False), encoding="utf-8")
    return str(target.resolve())


def publish(result: PreparationResult, path: str, **kwargs: Any) -> str:
    """Publish to whatever the file suffix asks for."""
    target = pathlib.Path(path)
    suffix = target.suffix.lower()

    if suffix == ".pdf":
        return to_pdf(result, path, **kwargs)
    if suffix == ".pptx":
        return to_pptx(result, path, **kwargs)
    if suffix == ".ipynb":
        return to_notebook(result, path, **kwargs)
    if suffix in (".html", ".htm", ".md", ".markdown"):
        return result.export_report(path)
    raise ValueError(
        f"cannot publish to {target.suffix!r}; expected .pdf, .pptx, .ipynb, .html or .md"
    )
